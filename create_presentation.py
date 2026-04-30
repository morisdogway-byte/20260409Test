from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
COLOR_BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy
COLOR_BG_MID  = RGBColor(0x16, 0x21, 0x3E)       # Mid navy
COLOR_ACCENT   = RGBColor(0xE9, 0x4F, 0x37)       # Red-orange
COLOR_ACCENT2  = RGBColor(0xF5, 0xA6, 0x23)       # Amber
COLOR_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT    = RGBColor(0xD0, 0xD8, 0xF0)       # Light blue-white
COLOR_MUTED    = RGBColor(0x90, 0xA0, 0xC0)       # Muted blue


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name="Noto Sans CJK JP"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color if color else COLOR_WHITE
    run.font.name = font_name
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=16, color=None, line_spacing=1.4,
                       font_name="Noto Sans CJK JP"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, opts) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", font_size))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", color if color else COLOR_WHITE)
        run.font.name = font_name
        if opts.get("space_before"):
            p.space_before = Pt(opts["space_before"])
    return txBox


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

# ──────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)

# Accent bar top
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)

# Large accent stripe
add_rect(s, 0, 2.8, 13.33, 2.1, COLOR_BG_MID)

# Side accent
add_rect(s, 0, 2.8, 0.15, 2.1, COLOR_ACCENT)

add_text_box(s, "ミニ講座  ／  30 min", 0.5, 0.35, 12, 0.5,
             font_size=14, color=COLOR_MUTED)

add_text_box(s, "チャールズ・デイビスの", 0.5, 1.1, 12, 0.8,
             font_size=30, bold=False, color=COLOR_LIGHT)

add_text_box(s, "明確さの原則", 0.5, 1.85, 12, 1.0,
             font_size=52, bold=True, color=COLOR_WHITE)

add_text_box(s, "Principle of Clarity", 0.5, 2.95, 12, 0.6,
             font_size=22, color=COLOR_ACCENT2)

add_text_box(s, "Charles Davis  ·  Bristol, UK\nhttps://medium.com/@charlesdavies",
             0.5, 5.5, 12, 0.9, font_size=14, color=COLOR_MUTED)

# Bottom bar
add_rect(s, 0, 7.3, 13.33, 0.2, COLOR_ACCENT)

# ──────────────────────────────────────────────
# SLIDE 2 — Agenda
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)

add_text_box(s, "本日の流れ", 0.5, 0.25, 12, 0.6,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.5, 0.95, 1.5, 0.05, COLOR_ACCENT)

items = [
    ("01", "チャールズ・デイビスとは？"),
    ("02", "ソース原理とは？"),
    ("03", "レスポンシビリティーの正体"),
    ("04", "明確さの原則とは？"),
    ("05", "痛みとモヤモヤの正体"),
    ("06", "2つの相反する行動"),
    ("07", "まとめ・実践へ"),
]

for i, (num, label) in enumerate(items):
    y = 1.2 + i * 0.76
    add_rect(s, 0.5, y, 0.55, 0.55, COLOR_ACCENT)
    add_text_box(s, num, 0.5, y + 0.05, 0.55, 0.45,
                 font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, label, 1.25, y + 0.1, 11, 0.45,
                 font_size=20, color=COLOR_LIGHT)

# ──────────────────────────────────────────────
# SLIDE 3 — Who is Charles Davis
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "01", 0.35, 0.2, 1.5, 0.55,
             font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "チャールズ・デイビスとは？", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 3, 0.05, COLOR_ACCENT2)

add_multiline_text(s, [
    ("👤  Charles Davis（チャールズ・デイビス）", {"size": 22, "bold": True, "color": COLOR_ACCENT2}),
    ("", {"size": 10}),
    ("🏙  イギリス・ブリストル在住", {"size": 20, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("📝  Medium :  https://medium.com/@charlesdavies", {"size": 16, "color": COLOR_MUTED}),
    ("", {"size": 10}),
    ("💡  「明確さの原則（Principle of Clarity）」の提唱者", {"size": 20, "color": COLOR_LIGHT}),
    ("", {"size": 10}),
    ("🤝  ソース原理の提唱者ピーター・カーニックとともに、\n    ソース原理を体系化したキーパーソン",
     {"size": 20, "color": COLOR_LIGHT}),
], 0.35, 1.55, 12.5, 5.5)

# ──────────────────────────────────────────────
# SLIDE 4 — Source Principle
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "02", 0.35, 0.2, 1.5, 0.55, font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "ソース原理とは？", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 2.5, 0.05, COLOR_ACCENT2)

# Box left
add_rect(s, 0.35, 1.6, 5.9, 4.8, COLOR_BG_MID)
add_rect(s, 0.35, 1.6, 0.12, 4.8, COLOR_ACCENT)

add_multiline_text(s, [
    ("ソース原理（Source Principle）", {"size": 19, "bold": True, "color": COLOR_ACCENT2}),
    ("", {"size": 8}),
    ("Peter Koenig（ピーター・カーニック）が\n提唱した原理。", {"size": 17, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("あらゆる取り組みには「ソース」と呼ばれる\n起点となる人物が存在し、その人が\nエネルギーと方向性を生み出す。", {"size": 17, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("チャールズはこの原理を\nさらに体系化する研究を行った。", {"size": 17, "color": COLOR_LIGHT}),
], 0.6, 1.75, 5.5, 4.5)

# Box right
add_rect(s, 6.6, 1.6, 6.3, 4.8, COLOR_BG_MID)
add_rect(s, 6.6, 1.6, 0.12, 4.8, COLOR_ACCENT2)

add_multiline_text(s, [
    ("チャールズの大きな発見", {"size": 19, "bold": True, "color": COLOR_ACCENT2}),
    ("", {"size": 8}),
    ("人が何かを始めようとする時、\nそこには必ず——", {"size": 17, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("「レスポンシビリティー」", {"size": 22, "bold": True, "color": COLOR_WHITE}),
    ("", {"size": 8}),
    ("——を感じている、という気づき。", {"size": 17, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("この発見がすべての出発点。", {"size": 17, "color": COLOR_MUTED}),
], 6.85, 1.75, 5.9, 4.5)

# ──────────────────────────────────────────────
# SLIDE 5 — What is Responsibility
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "03", 0.35, 0.2, 1.5, 0.55, font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "レスポンシビリティーの正体", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 3.5, 0.05, COLOR_ACCENT2)

# NOT simple accountability
add_rect(s, 0.35, 1.55, 12.6, 1.1, RGBColor(0x3A, 0x10, 0x10))
add_rect(s, 0.35, 1.55, 0.12, 1.1, COLOR_ACCENT)
add_multiline_text(s, [
    ("❌  単なる「説明責任（accountability）」ではない！", {"size": 18, "bold": True, "color": COLOR_ACCENT}),
    ("    レスポンシビリティーとは——筋肉が「レスポンス（反応）」してしまう、身体的な感覚のこと。",
     {"size": 17, "color": COLOR_LIGHT}),
], 0.5, 1.65, 12.3, 1.0)

# Example box
add_rect(s, 0.35, 2.9, 12.6, 3.5, COLOR_BG_MID)
add_rect(s, 0.35, 2.9, 0.12, 3.5, COLOR_ACCENT2)

add_multiline_text(s, [
    ("🔥  体験ワーク  ——  想像してみてください", {"size": 19, "bold": True, "color": COLOR_ACCENT2}),
    ("", {"size": 8}),
    ("キッチンで火をつけて料理をしている最中に……", {"size": 18, "color": COLOR_LIGHT}),
    ("", {"size": 6}),
    ("「火はつけっぱなしのまま、今すぐ電車に乗って出かけてください」", {"size": 18, "bold": True, "color": COLOR_WHITE}),
    ("", {"size": 8}),
    ("そう言われたら、どんな身体感覚がありますか？", {"size": 18, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("👉  ソワソワ、落ち着かない……  その筋肉レベルの感覚こそが", {"size": 17, "color": COLOR_LIGHT}),
    ("    「レスポンシビリティー」の正体です。", {"size": 20, "bold": True, "color": COLOR_ACCENT2}),
], 0.6, 3.05, 12.3, 3.3)

# ──────────────────────────────────────────────
# SLIDE 6 — Principle of Clarity intro
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "04", 0.35, 0.2, 1.5, 0.55, font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "明確さの原則とは？", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 2.8, 0.05, COLOR_ACCENT2)

add_multiline_text(s, [
    ("レスポンシビリティーを感じていても……", {"size": 20, "color": COLOR_MUTED}),
    ("", {"size": 6}),
    ("「何に対して」レスポンシビリティーを感じているのかが\n分からなければ、意味がない。", {"size": 22, "color": COLOR_LIGHT}),
], 0.35, 1.55, 12.6, 1.8)

# Arrow
add_rect(s, 6.3, 3.5, 0.7, 0.5, COLOR_ACCENT)
add_text_box(s, "▼", 6.3, 3.5, 0.7, 0.5,
             font_size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# Solution box
add_rect(s, 0.35, 4.1, 12.6, 2.7, COLOR_BG_MID)
add_rect(s, 0.35, 4.1, 0.12, 2.7, COLOR_ACCENT2)

add_multiline_text(s, [
    ("✨  明確さの原則（Principle of Clarity）", {"size": 22, "bold": True, "color": COLOR_ACCENT2}),
    ("", {"size": 8}),
    ("チャールズが発明した「レスポンシビリティーを明確にするためのワーク」。", {"size": 19, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("自分が何に反応しているのかを言語化・可視化し、\nエネルギーを正しい方向へ向けるためのフレームワーク。",
     {"size": 19, "color": COLOR_LIGHT}),
    ("", {"size": 8}),
    ("詳細 ▶  https://note.com/aonohideaki/n/n3310f5e50256",
     {"size": 14, "color": COLOR_MUTED}),
], 0.6, 4.25, 12.3, 2.5)

# ──────────────────────────────────────────────
# SLIDE 7 — Pain & Moya-Moya
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "05", 0.35, 0.2, 1.5, 0.55, font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "痛みとモヤモヤの正体", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 2.8, 0.05, COLOR_ACCENT2)

add_multiline_text(s, [
    ("チャールズのさらなる発見——", {"size": 20, "color": COLOR_MUTED}),
    ("", {"size": 6}),
    ("レスポンシビリティーが発揮されていない時、\n人は「痛み（モヤモヤ・不快感）」を感じている。",
     {"size": 22, "bold": True, "color": COLOR_WHITE}),
], 0.35, 1.55, 12.6, 1.7)

# Highlight
add_rect(s, 0.35, 3.45, 12.6, 1.0, RGBColor(0x2A, 0x1A, 0x00))
add_rect(s, 0.35, 3.45, 0.12, 1.0, COLOR_ACCENT2)
add_text_box(s, "「やりたいのにやらない」「やりたくないのにやっている」\nこの2つが、痛みとモヤモヤを生じさせる源泉。",
             0.55, 3.55, 12.3, 0.85, font_size=19, bold=True, color=COLOR_ACCENT2)

add_text_box(s, "つまり——レスポンシビリティーと行動が一致しない状態が「痛み」の正体。",
             0.35, 4.65, 12.6, 0.6, font_size=18, color=COLOR_LIGHT)

add_rect(s, 0.35, 5.45, 12.6, 1.7, COLOR_BG_MID)
add_text_box(s, "💡  逆に言えば、明確さの原則でレスポンシビリティーを可視化し、\n    行動と一致させることができれば——痛みは消え、力が生まれる。",
             0.55, 5.6, 12.2, 1.4, font_size=18, color=COLOR_LIGHT)

# ──────────────────────────────────────────────
# SLIDE 8 — Two contradictory behaviors
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "06", 0.35, 0.2, 1.5, 0.55, font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "2つの相反する行動", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 2.5, 0.05, COLOR_ACCENT2)

# Pattern A
add_rect(s, 0.35, 1.6, 5.9, 4.5, RGBColor(0x2A, 0x10, 0x10))
add_rect(s, 0.35, 1.6, 0.12, 4.5, COLOR_ACCENT)
add_multiline_text(s, [
    ("パターン A", {"size": 16, "bold": True, "color": COLOR_ACCENT}),
    ("", {"size": 8}),
    ("やりたいと思っているのに、", {"size": 20, "color": COLOR_LIGHT}),
    ("やらない", {"size": 28, "bold": True, "color": COLOR_WHITE}),
    ("", {"size": 10}),
    ("例）副業を始めたいと思いながら\n　　何ヶ月も先延ばしにしている……", {"size": 17, "color": COLOR_MUTED}),
    ("", {"size": 10}),
    ("→ 内側では「やりたい」という\n　 レスポンシビリティーが動いているのに\n　 行動が伴っていない状態。",
     {"size": 16, "color": COLOR_LIGHT}),
], 0.6, 1.75, 5.5, 4.3)

# Pattern B
add_rect(s, 6.6, 1.6, 6.3, 4.5, RGBColor(0x10, 0x10, 0x2A))
add_rect(s, 6.6, 1.6, 0.12, 4.5, COLOR_ACCENT2)
add_multiline_text(s, [
    ("パターン B", {"size": 16, "bold": True, "color": COLOR_ACCENT2}),
    ("", {"size": 8}),
    ("やりたくないと思っているのに、", {"size": 20, "color": COLOR_LIGHT}),
    ("やっている", {"size": 28, "bold": True, "color": COLOR_WHITE}),
    ("", {"size": 10}),
    ("例）本当はNOと言いたい依頼を\n　　断れずに引き受け続けている……", {"size": 17, "color": COLOR_MUTED}),
    ("", {"size": 10}),
    ("→ レスポンシビリティーがないことを\n　 やり続けることで蓄積される\n　 消耗と痛み。",
     {"size": 16, "color": COLOR_LIGHT}),
], 6.85, 1.75, 5.9, 4.3)

# Bottom summary
add_rect(s, 0.35, 6.3, 12.6, 0.95, COLOR_BG_MID)
add_text_box(s, "この2つは「レスポンシビリティーと行動のズレ」という共通の根を持つ。\n明確さの原則は、このズレを解消するための鍵。",
             0.55, 6.4, 12.2, 0.8, font_size=17, color=COLOR_LIGHT)

# ──────────────────────────────────────────────
# SLIDE 9 — Summary / Call to Action
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 0.08, 0.15, 7.42, COLOR_ACCENT)

add_text_box(s, "07", 0.35, 0.2, 1.5, 0.55, font_size=14, color=COLOR_ACCENT2)
add_text_box(s, "まとめ・実践へ", 0.35, 0.6, 12, 0.7,
             font_size=32, bold=True, color=COLOR_WHITE)
add_rect(s, 0.35, 1.35, 2.2, 0.05, COLOR_ACCENT2)

summary = [
    ("レスポンシビリティーは「説明責任」ではなく筋肉レベルの身体感覚", COLOR_LIGHT),
    ("何かを始めようとする時、人は必ずレスポンシビリティーを感じている", COLOR_LIGHT),
    ("「何に対して」感じているかを明確にする必要がある", COLOR_LIGHT),
    ("明確さの原則はそのための実践的なワーク", COLOR_ACCENT2),
    ("行動とレスポンシビリティーのズレが痛みとモヤモヤを生む", COLOR_LIGHT),
    ("明確にすることで、エネルギーが解放される", COLOR_ACCENT2),
]

for i, (text, col) in enumerate(summary):
    y = 1.6 + i * 0.72
    add_rect(s, 0.35, y, 0.45, 0.45, COLOR_ACCENT if col == COLOR_LIGHT else COLOR_ACCENT2)
    add_text_box(s, str(i+1), 0.35, y + 0.05, 0.45, 0.38,
                 font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, text, 0.95, y + 0.08, 12, 0.38,
                 font_size=18, color=col)

# Reference
add_text_box(s, "📖  詳しくは ▶  https://note.com/aonohideaki/n/n3310f5e50256",
             0.35, 6.95, 12.6, 0.45, font_size=14, color=COLOR_MUTED)

# ──────────────────────────────────────────────
# SLIDE 10 — Closing
# ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, COLOR_BG_DARK)
add_rect(s, 0, 0, 13.33, 0.08, COLOR_ACCENT)
add_rect(s, 0, 3.0, 13.33, 0.08, COLOR_ACCENT2)

add_text_box(s, "Thank You", 0, 1.0, 13.33, 1.2,
             font_size=60, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

add_text_box(s, "ご清聴ありがとうございました", 0, 2.2, 13.33, 0.7,
             font_size=26, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

add_multiline_text(s, [
    ("Charles Davis  |  Principle of Clarity", {"size": 18, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}),
    ("", {"size": 10}),
    ("https://medium.com/@charlesdavies", {"size": 15, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}),
    ("https://note.com/aonohideaki/n/n3310f5e50256", {"size": 15, "color": COLOR_MUTED, "align": PP_ALIGN.CENTER}),
], 0, 3.5, 13.33, 3.5, font_size=18)

add_rect(s, 0, 7.3, 13.33, 0.2, COLOR_ACCENT)

# ── Save ──
output_path = "/home/user/20260409Test/チャールズ・デイビスの明確さの原則.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
